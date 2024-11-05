import tkinter as tk
from docx2pdf import convert as docx_to_pdf
import comtypes.client
import zipfile
import fitz
from pptx.util import Inches
from pptx import Presentation
from pptx.util import Pt
from PIL import Image
import tempfile
import os
from tkinter import filedialog, messagebox

# Fungsi Gabung JPG dan Konversi ke PDF
def jpg_to_pdf():
    files = filedialog.askopenfilenames(filetypes=[("JPG Files", "*.jpg")])
    if not files:
        return
    images = [Image.open(f).convert("RGB") for f in files]
    pdf_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
    if pdf_path:
        images[0].save(pdf_path, save_all=True, append_images=images[1:])
        messagebox.showinfo("Sukses", "JPG berhasil dikonversi ke PDF!")

# Fungsi Konversi Word ke PDF
def word_to_pdf():
    file = filedialog.askopenfilename(filetypes=[("Word Files", "*.docx")])
    if file:
        pdf_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
        if pdf_path:
            docx_to_pdf(file, pdf_path)
            messagebox.showinfo("Sukses", "Word berhasil dikonversi ke PDF!")

# Fungsi Konversi PPT ke PDF
def ppt_to_pdf():
    file = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
    if file:
        pdf_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
        if pdf_path:
            # Gunakan path absolut untuk menghindari masalah path
            file = os.path.abspath(file)
            pdf_path = os.path.abspath(pdf_path)

            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            try:
                ppt = powerpoint.Presentations.Open(file)
                ppt.SaveAs(pdf_path, 32)  # 32 = pdf format
                ppt.Close()
                messagebox.showinfo("Sukses", "PPT berhasil dikonversi ke PDF!")
            except Exception as e:
                messagebox.showerror("Error", f"Gagal mengonversi: {e}")
            finally:
                powerpoint.Quit()

# Fungsi Kompres PDF
def compress_pdf():
    file = filedialog.askopenfilename(filetypes=[("PDF Files", "*.pdf")])
    if file:
        pdf_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
        if pdf_path:
            pdf_document = fitz.open(file)
            
            # Simpan ke PDF baru dengan kompresi gambar dan penghapusan konten sampah
            pdf_document.save(pdf_path, garbage=4, deflate=True, deflate_images=True)
            pdf_document.close()
            
            messagebox.showinfo("Sukses", "PDF berhasil dikompres dengan ukuran lebih kecil!")

# Fungsi Kompres PPT
def compress_ppt():
    file = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
    if file:
        ppt = Presentation(file)
        output_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")])

        if output_path:
            for slide in ppt.slides:
                for shape in slide.shapes:
                    # Kompres gambar
                    if shape.shape_type == 13:  # Bentuk gambar
                        image = shape.image
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp_img:
                            # Simpan gambar sementara dengan kualitas rendah
                            image_bytes = image.blob
                            with open(tmp_img.name, "wb") as f:
                                f.write(image_bytes)

                            # Buka gambar dan kompres
                            img = Image.open(tmp_img.name)
                            if img.mode == 'RGBA':  # Konversi dari RGBA ke RGB jika perlu
                                img = img.convert('RGB')
                            img = img.resize((int(img.width * 0.4), int(img.height * 0.4)), Image.LANCZOS)  # Resize ke 40%
                            img.save(tmp_img.name, quality=50)  # Simpan dengan kualitas gambar lebih rendah

                            # Gantikan gambar lama dengan yang baru
                            shape.element.getparent().remove(shape.element)
                            slide.shapes.add_picture(tmp_img.name, shape.left, shape.top, shape.width, shape.height)

                        os.remove(tmp_img.name)  # Hapus file sementara

                    # Kurangi ukuran teks jika ada
                    elif shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size and run.font.size.pt > 10:  # Atur batas minimum untuk ukuran font
                                    run.font.size = Pt(max(10, run.font.size.pt * 0.8))  # Kurangi ukuran font sebesar 80%
                    
                    # Hapus elemen bentuk lain yang tidak perlu (opsional)
                    if shape.shape_type not in [1, 13]:  # Misal, hapus elemen selain teks dan gambar
                        shape._element.getparent().remove(shape._element)

            # Simpan hasilnya
            ppt.save(output_path)
            messagebox.showinfo("Sukses", "PPT berhasil dikompres dengan ukuran lebih kecil!")

# GUI dengan Tkinter
def create_gui():
    window = tk.Tk()
    window.title("File Converter & Compressor")
    window.geometry("400x400")

    tk.Button(window, text="Gabung JPG & Konversi ke PDF", command=jpg_to_pdf).pack(pady=10)
    tk.Button(window, text="Konversi Word ke PDF", command=word_to_pdf).pack(pady=10)
    tk.Button(window, text="Konversi PPT ke PDF", command=ppt_to_pdf).pack(pady=10)
    tk.Button(window, text="Kompres PDF", command=compress_pdf).pack(pady=10)
    tk.Button(window, text="Kompres PPT", command=compress_ppt).pack(pady=10)

    window.mainloop()

if __name__ == "__main__":
    create_gui()